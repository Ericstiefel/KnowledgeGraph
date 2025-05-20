from pptx import Presentation
from PIL import Image
from io import BytesIO

def extract_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    ordered_content = []

    for slide_number, slide in enumerate(prs.slides):
        slide_items = []

        for shape in slide.shapes:
            top = shape.top if hasattr(shape, 'top') else 0

            # Extract text
            if shape.has_text_frame:
                text = "\n".join([p.text for p in shape.text_frame.paragraphs if p.text])
                if text.strip():
                    slide_items.append(("text", top, text))

            # Extract image
            elif shape.shape_type == 13 and hasattr(shape, "image"):  # Picture shape
                image = shape.image
                image_bytes = BytesIO(image.blob)
                img = Image.open(image_bytes)

                slide_items.append(("image", top, img))  # Store PIL image directly

        # Sort items by vertical position
        slide_items.sort(key=lambda x: x[1])

        for kind, _, content in slide_items:
            ordered_content.append((kind, content))

    return ordered_content
